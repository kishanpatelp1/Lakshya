"""Document processing service -- parse, chunk, embed, and index uploaded files."""

import csv
import io
import logging
import uuid
from pathlib import Path
from typing import Any

from src.config import get_settings
from src.llm import get_embeddings, get_embedding_dim

logger = logging.getLogger(__name__)

CHUNK_SIZE = 512
CHUNK_OVERLAP = 64


class DocumentProcessor:
    """Parse uploaded documents, chunk text, embed, and upsert into Qdrant."""

    COLLECTION = "user_uploads"

    def __init__(self) -> None:
        self.settings = get_settings()
        self._qdrant = None

    def _get_qdrant(self):
        if self._qdrant is None:
            from qdrant_client import QdrantClient

            self._qdrant = QdrantClient(
                url=self.settings.qdrant_url,
                api_key=self.settings.qdrant_api_key,
            )
        return self._qdrant

    def _ensure_collection(self) -> None:
        client = self._get_qdrant()
        from qdrant_client.models import Distance, VectorParams

        collections = [c.name for c in client.get_collections().collections]
        if self.COLLECTION not in collections:
            client.create_collection(
                collection_name=self.COLLECTION,
                vectors_config=VectorParams(
                    size=get_embedding_dim(),
                    distance=Distance.COSINE,
                ),
            )

    # ── public entry point ──────────────────────────────────────────────

    def process(
        self,
        file_path: str,
        user_id: str,
        upload_id: str,
        file_type: str,
    ) -> int:
        """Parse *file_path*, chunk, embed, and upsert.  Returns chunk count."""
        pages = self._extract_pages(file_path, file_type)
        if not pages:
            logger.warning("No text extracted from %s", file_path)
            return 0

        chunks = self._chunk_pages(pages)
        if not chunks:
            return 0

        self._ensure_collection()
        self._upsert_chunks(chunks, user_id, upload_id)
        return len(chunks)

    # ── extraction ──────────────────────────────────────────────────────

    def _extract_pages(
        self, file_path: str, file_type: str
    ) -> list[dict[str, Any]]:
        """Return list of {page_number, title, text}."""
        ft = file_type.lower().lstrip(".")
        if ft == "pdf":
            return self._parse_pdf(file_path)
        if ft in ("pptx", "ppt"):
            return self._parse_pptx(file_path)
        if ft == "txt":
            return self._parse_txt(file_path)
        if ft == "csv":
            return self._parse_csv(file_path)
        if ft == "xlsx":
            return self._parse_xlsx(file_path)
        logger.warning("Unsupported file type: %s", ft)
        return []

    @staticmethod
    def _parse_pdf(path: str) -> list[dict[str, Any]]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.error("PyMuPDF (fitz) not installed – pip install PyMuPDF")
            return []

        pages: list[dict[str, Any]] = []
        with fitz.open(path) as doc:
            for idx, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if text:
                    pages.append({"page_number": idx, "title": f"Page {idx}", "text": text})
        return pages

    @staticmethod
    def _parse_pptx(path: str) -> list[dict[str, Any]]:
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed – pip install python-pptx")
            return []

        pages: list[dict[str, Any]] = []
        prs = Presentation(path)
        for idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    frame_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                    if frame_text:
                        parts.append(frame_text)
                if hasattr(shape, "name") and "title" in shape.name.lower() and shape.has_text_frame:
                    slide_title = shape.text_frame.text.strip()
            text = "\n".join(parts).strip()
            if text:
                pages.append({
                    "page_number": idx,
                    "title": slide_title or f"Slide {idx}",
                    "text": text,
                })
        return pages

    @staticmethod
    def _parse_txt(path: str) -> list[dict[str, Any]]:
        text = Path(path).read_text(encoding="utf-8", errors="replace").strip()
        if not text:
            return []
        return [{"page_number": 1, "title": "Document", "text": text}]

    @staticmethod
    def _parse_csv(path: str) -> list[dict[str, Any]]:
        raw = Path(path).read_text(encoding="utf-8", errors="replace")
        reader = csv.reader(io.StringIO(raw))
        lines = [", ".join(row) for row in reader if any(cell.strip() for cell in row)]
        text = "\n".join(lines).strip()
        if not text:
            return []
        return [{"page_number": 1, "title": "CSV Data", "text": text}]

    @staticmethod
    def _parse_xlsx(path: str) -> list[dict[str, Any]]:
        try:
            import openpyxl
        except ImportError:
            logger.error("openpyxl not installed – pip install openpyxl")
            return []

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        pages: list[dict[str, Any]] = []
        for sheet_idx, ws in enumerate(wb.worksheets, start=1):
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(", ".join(cells))
            text = "\n".join(rows).strip()
            if text:
                pages.append({
                    "page_number": sheet_idx,
                    "title": ws.title or f"Sheet {sheet_idx}",
                    "text": text,
                })
        wb.close()
        return pages

    # ── chunking ────────────────────────────────────────────────────────

    @staticmethod
    def _chunk_pages(
        pages: list[dict[str, Any]],
    ) -> list[dict[str, Any]]:
        """Split each page's text into overlapping chunks, preserving metadata."""
        chunks: list[dict[str, Any]] = []
        chunk_idx = 0
        for page in pages:
            words = page["text"].split()
            if not words:
                continue
            start = 0
            while start < len(words):
                end = start + CHUNK_SIZE
                chunk_text = " ".join(words[start:end])
                chunks.append({
                    "page_number": page["page_number"],
                    "title": page["title"],
                    "text": chunk_text,
                    "chunk_index": chunk_idx,
                })
                chunk_idx += 1
                step = CHUNK_SIZE - CHUNK_OVERLAP
                if step <= 0:
                    step = CHUNK_SIZE
                start += step
        return chunks

    # ── embedding + upsert ──────────────────────────────────────────────

    def _upsert_chunks(
        self,
        chunks: list[dict[str, Any]],
        user_id: str,
        upload_id: str,
    ) -> None:
        embeddings_model = get_embeddings()
        texts = [c["text"] for c in chunks]

        batch_size = 32
        all_vectors: list[list[float]] = []
        for i in range(0, len(texts), batch_size):
            batch = texts[i : i + batch_size]
            all_vectors.extend(embeddings_model.embed_documents(batch))

        from qdrant_client.models import PointStruct

        points = []
        for chunk, vector in zip(chunks, all_vectors):
            points.append(
                PointStruct(
                    id=str(uuid.uuid4()),
                    vector=vector,
                    payload={
                        "user_id": user_id,
                        "upload_id": upload_id,
                        "text": chunk["text"],
                        "page_number": chunk["page_number"],
                        "slide_title": chunk["title"],
                        "chunk_index": chunk["chunk_index"],
                    },
                )
            )

        client = self._get_qdrant()
        upsert_batch = 64
        for i in range(0, len(points), upsert_batch):
            client.upsert(
                collection_name=self.COLLECTION,
                points=points[i : i + upsert_batch],
            )
        logger.info(
            "Indexed %d chunks for upload %s (user %s)",
            len(points),
            upload_id,
            user_id,
        )
