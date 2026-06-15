"""Document parsing tools for PDFs, PPTs, and URLs."""

from pathlib import Path
from typing import Any, Dict, List

from langchain_core.tools import tool


@tool
def parse_pdf(pdf_path: str) -> Dict[str, Any]:
    """Parse a PDF file and extract text, tables, and numeric data.

    Args:
        pdf_path: Absolute path to the PDF file on disk
    """
    import re

    p = Path(pdf_path)
    if not p.exists() or not p.is_file():
        return {"error": f"PDF not found: {pdf_path}"}

    try:
        import fitz  # type: ignore

        doc = fitz.open(str(p))
        pages: List[Dict[str, Any]] = []
        numbers: List[float] = []
        for i, page in enumerate(doc):
            text = page.get_text("text") or ""
            pages.append({"page": i + 1, "text": text})
            for tok in re.findall(r"[-+]?\d*\.?\d+", text):
                try:
                    numbers.append(float(tok))
                except ValueError:
                    continue

        tables: List[Dict[str, Any]] = []
        try:
            import pdfplumber  # type: ignore

            with pdfplumber.open(str(p)) as pdf:
                for pg_idx, pg in enumerate(pdf.pages, start=1):
                    for table in pg.extract_tables() or []:
                        if table:
                            tables.append({"page": pg_idx, "rows": table})
        except Exception:
            pass

        return {
            "source": str(p),
            "doc_type": "pdf",
            "text_pages": pages,
            "tables": tables,
            "numbers_count": len(numbers),
        }
    except Exception as exc:
        return {"error": str(exc), "source": str(p)}


@tool
def parse_ppt(ppt_path: str) -> Dict[str, Any]:
    """Parse a PowerPoint file and extract slide text and numeric data.

    Args:
        ppt_path: Absolute path to the PPT/PPTX file on disk
    """
    import re

    p = Path(ppt_path)
    if not p.exists() or not p.is_file():
        return {"error": f"PPT not found: {ppt_path}"}

    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(p))
        slides: List[Dict[str, Any]] = []
        numbers: List[float] = []
        for idx, slide in enumerate(prs.slides, start=1):
            chunks = []
            for shape in slide.shapes:
                txt = getattr(shape, "text", "")
                if txt:
                    chunks.append(txt)
                    for tok in re.findall(r"[-+]?\d*\.?\d+", txt):
                        try:
                            numbers.append(float(tok))
                        except ValueError:
                            continue
            slides.append({"slide": idx, "text": "\n".join(chunks)})

        return {
            "source": str(p),
            "doc_type": "ppt",
            "slides": slides,
            "numbers_count": len(numbers),
        }
    except Exception as exc:
        return {"error": str(exc), "source": str(p)}


@tool
def fetch_url(url: str) -> Dict[str, Any]:
    """Fetch and extract text content from a URL.

    Args:
        url: The URL to fetch
    """
    import re

    import requests

    try:
        resp = requests.get(url, timeout=25)
        resp.raise_for_status()
        text = ""
        try:
            from bs4 import BeautifulSoup  # type: ignore

            soup = BeautifulSoup(resp.text, "html.parser")
            text = " ".join(soup.stripped_strings)
        except ImportError:
            text = resp.text

        numbers_count = len(re.findall(r"[-+]?\d*\.?\d+", text))
        return {
            "source": url,
            "doc_type": "url",
            "text": text[:100_000],
            "char_count": len(text),
            "numbers_count": numbers_count,
        }
    except Exception as exc:
        return {"error": str(exc), "source": url}
