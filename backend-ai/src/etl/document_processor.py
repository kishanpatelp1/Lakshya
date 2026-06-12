"""Document processing pipeline implementation for the ETL Transform layer."""

import hashlib
import logging
import re
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
import pandas as pd
from pptx import Presentation

# Optional deps — PDF extraction uses pdfplumber (see extract_text_from_pdf),
# so these are only needed for their specific formats. Guard the imports so a
# missing PyMuPDF/PyPDF2/python-docx doesn't break the whole ETL pipeline.
try:
    import fitz  # PyMuPDF (unused; kept for optional callers)
except ImportError:  # pragma: no cover
    fitz = None
try:
    from PyPDF2 import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None
try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Main document processing class for handling different document types."""
    
    def __init__(self):
        self.processed_dir = Path("uploads/filings")
        self.processed_dir.mkdir(exist_ok=True)

    def process_document(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Process a document and return extracted content."""
        try:
            path = Path(file_path)
            if not path.exists():
                logger.error(f"Document not found: {file_path}")
                return None
                
            # Determine file type and process accordingly
            if path.suffix.lower() == '.pdf':
                text = self.extract_text_from_pdf(str(path))
                return {"text": text, "file_type": "pdf"}
            elif path.suffix.lower() == '.docx':
                text = self.extract_text_from_docx(str(path))
                return {"text": text, "file_type": "docx"}
            elif path.suffix.lower() == '.pptx':
                text = self.extract_text_from_pptx(str(path))
                return {"text": text, "file_type": "pptx"}
            elif path.suffix.lower() in ['.txt', '.text']:
                text = self.parse_text(str(path))
                return {"text": text, "file_type": "text"}
            else:
                logger.warning(f"Unsupported file type: {path.suffix}")
                return None
                
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            return None

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text and tables from PDF using pdfplumber."""
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    
                    tables = page.extract_tables()
                    for table in tables:
                        if not table:
                            continue
                        text += "\n"
                        for i, row in enumerate(table):
                            clean_row = [str(cell).replace('\n', ' ') if cell else '' for cell in row]
                            text += "| " + " | ".join(clean_row) + " |\n"
                            if i == 0:  # Add markdown table header separator
                                text += "|" + "|".join(["---"] * len(clean_row)) + "|\n"
                        text += "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return ""

    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            return ""

    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
            return ""

    def parse_text(self, file_path: str) -> str:
        """Parse text file and return content."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return ""